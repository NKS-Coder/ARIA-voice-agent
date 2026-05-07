#!/usr/bin/env python3
"""Builds ARIA-Demo.pptx — a deck covering the project evolution, architecture,
and tools. Designed for a 10-12 minute professor demo."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x04, 0x05, 0x0A)
TEAL = RGBColor(0x5D, 0xE6, 0xD8)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
WHITE = RGBColor(0xF0, 0xF2, 0xFC)
DIM = RGBColor(0x6B, 0x7E, 0xAA)
ACCENT = RGBColor(0xFF, 0xD7, 0x00)
RED = RGBColor(0xFC, 0x81, 0x81)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_bg(slide, color=NAVY):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    if "\n" in text:
        lines = text.split("\n")
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run(); run.text = lines[0]
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
        for line in lines[1:]:
            p = tf.add_paragraph(); p.alignment = align
            run = p.add_run(); run.text = line
            run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
    else:
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=WHITE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run(); run.text = "• " + item
        run.font.size = Pt(size); run.font.color.rgb = color; run.font.name = "Calibri"


def add_card(slide, x, y, w, h, fill=RGBColor(0x10, 0x14, 0x22), line=RGBColor(0x2A, 0x33, 0x4D)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.05
    card.fill.solid(); card.fill.fore_color.rgb = fill
    card.line.color.rgb = line; card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def slide_header(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.5), Inches(0.6),
             title, size=30, bold=True, color=TEAL)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.4),
                 subtitle, size=14, color=DIM)
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.45), Inches(12.8), Inches(1.45))
    line.line.color.rgb = TEAL; line.line.width = Pt(1.5)


# ── Slide 1 — Title ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
glow = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(1.7), Inches(3), Inches(3))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x10, 0x1F, 0x3A); glow.line.fill.background()
add_text(s, Inches(0), Inches(2.2), Inches(13.33), Inches(1.2),
         "ARIA", size=88, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.5), Inches(13.33), Inches(0.5),
         "Neural Command Interface", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(4.1), Inches(13.33), Inches(0.4),
         "A voice-first AI assistant for Gmail, Calendar, and beyond",
         size=16, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.5), Inches(13.33), Inches(0.4),
         "Nihhar K Sonee  ·  May 2026", size=14, color=TEAL, align=PP_ALIGN.CENTER)

# ── Slide 2 — The problem ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "The problem", "Inboxes are overwhelming. Voice should fix that.")
add_card(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(2.4))
add_text(s, Inches(1), Inches(2.1), Inches(5.4), Inches(0.5),
         "What users actually want", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(1), Inches(2.7), Inches(5.5), Inches(1.6), [
    "\"What's important in my inbox today?\"",
    "\"Send a quick reply to John\"",
    "\"Archive everything from LinkedIn\"",
    "\"Show me my bank statements\"",
], size=14)
add_card(s, Inches(6.8), Inches(1.9), Inches(5.9), Inches(2.4))
add_text(s, Inches(7.1), Inches(2.1), Inches(5.4), Inches(0.5),
         "What they're stuck doing", size=18, bold=True, color=RED)
add_bullets(s, Inches(7.1), Inches(2.7), Inches(5.5), Inches(1.6), [
    "Scrolling through 200 unread emails",
    "Tapping on tiny mobile UIs",
    "Switching apps to send a one-line reply",
    "Manually filtering financial vs. job vs. spam",
], size=14)
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.6),
         "ARIA is the voice layer between you and your inbox.",
         size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(5.4), Inches(12), Inches(0.6),
         "Speak naturally. Get filtered, ranked, summarized results.",
         size=16, color=DIM, align=PP_ALIGN.CENTER)

# ── Slide 3 — Architecture ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Architecture", "Browser ⇄ Cloudflare Worker ⇄ Groq + ElevenLabs + Google APIs")
# 3 columns
def arch_col(x, label, items, color):
    add_card(s, Inches(x), Inches(1.9), Inches(4.0), Inches(4.6))
    add_text(s, Inches(x+0.2), Inches(2.05), Inches(3.7), Inches(0.5),
             label, size=18, bold=True, color=color)
    add_bullets(s, Inches(x+0.2), Inches(2.7), Inches(3.7), Inches(3.7),
                items, size=13)
arch_col(0.5, "Frontend (browser)", [
    "Static HTML/CSS/JS — no framework",
    "Web Speech API (free, browser-native)",
    "Gmail OAuth (PKCE flow)",
    "GitHub Pages hosting",
    "Single-recognizer wake word",
    "Editable confirm-send modal",
], TEAL)
arch_col(4.7, "Cloudflare Worker", [
    "Routing brain (detectIntent regex)",
    "Importance ranker (7 parallel queries)",
    "Email composer (Gmail API)",
    "TTS proxy (ElevenLabs)",
    "OAuth callback handler",
    "Supabase token persistence",
], CYAN)
arch_col(8.9, "External services", [
    "Groq (Llama 3.3 70B, fallback 8B)",
    "ElevenLabs (Bella voice)",
    "Gmail API (search, send, organize)",
    "Google OAuth 2.0",
    "Supabase (per-session tokens)",
    "Cloudflare Workers (edge runtime)",
], ACCENT)

# ── Slide 4 — The voice loop ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "The voice loop", "From wake word to action in under 3 seconds")
steps = [
    ("1", "Wake", '\"Hey ARIA\"', "Web Speech API"),
    ("2", "Listen", "10s command window", "Browser STT"),
    ("3", "Route", "detectIntent(msg)", "Cloudflare Worker"),
    ("4", "Search/Act", "Gmail API call", "Worker → Google"),
    ("5", "Speak", "Stream MP3 reply", "ElevenLabs TTS"),
]
y = Inches(2.2); xw = Inches(2.5); gap = Inches(0.05); total = len(steps) * (xw + gap) - gap
start_x = (prs.slide_width - total) / 2
for i, (num, head, sub, tech) in enumerate(steps):
    x = start_x + i * (xw + gap)
    card = add_card(s, x, y, xw, Inches(2.6))
    add_text(s, x + Inches(0.2), y + Inches(0.15), xw - Inches(0.4), Inches(0.6),
             num, size=44, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(0.95), xw - Inches(0.2), Inches(0.4),
             head, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(1.4), xw - Inches(0.2), Inches(0.5),
             sub, size=12, color=DIM, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.1), y + Inches(2.05), xw - Inches(0.2), Inches(0.4),
             tech, size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5),
         "Latency budget: ~280ms STT + ~700ms Groq + ~1.2s TTS = ~2.2s end-to-end",
         size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# ── Slide 5 — Tools and why ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Tools — and why", "Every choice optimizes for cost, latency, or developer ergonomics")

tools = [
    ("Groq (Llama 3.3 70B)", "LLM inference",
     "10x faster than OpenAI at ~1/10 the cost. Llama 3.3 is open-weights, no vendor lock-in. Auto-fallback to 3.1-8b on rate limits."),
    ("Cloudflare Workers", "Backend runtime",
     "Zero cold-start, edge-deployed (low TTFB anywhere), 100k free requests/day. Routes through one /chat endpoint."),
    ("ElevenLabs", "Text-to-speech",
     "Highest-quality voice on the market. Bella voice for warmth. Streaming MP3 = TTS starts playing in ~1s."),
    ("Web Speech API", "Speech-to-text",
     "Free, browser-native, no API key. Continuous recognition + interim results. Wake word runs entirely client-side."),
    ("Gmail API + OAuth 2.0", "Email backend",
     "Direct, no third-party broker. PKCE flow keeps tokens browser-side; Supabase persists refresh tokens."),
    ("GitHub Pages", "Frontend hosting",
     "Free, fast CDN, automatic HTTPS. Static HTML keeps the surface area minimal."),
]
y = Inches(1.7)
for i, (tool, role, why) in enumerate(tools):
    row = i // 2; col = i % 2
    cy = y + Inches(row * 1.85)
    cx = Inches(0.5 + col * 6.4)
    add_card(s, cx, cy, Inches(6.2), Inches(1.7))
    add_text(s, cx + Inches(0.2), cy + Inches(0.1), Inches(5.9), Inches(0.4),
             tool, size=15, bold=True, color=TEAL)
    add_text(s, cx + Inches(0.2), cy + Inches(0.5), Inches(5.9), Inches(0.35),
             role, size=11, color=ACCENT)
    add_text(s, cx + Inches(0.2), cy + Inches(0.9), Inches(5.9), Inches(0.75),
             why, size=11, color=WHITE)

# ── Slide 6 — Smart importance ranker ──────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Smart importance ranker", "From \"top 5 emails\" to \"top 5 BANK emails of this month\"")
add_card(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.0))
add_text(s, Inches(0.7), Inches(1.85), Inches(5.6), Inches(0.5),
         "How it works", size=18, bold=True, color=TEAL)
add_bullets(s, Inches(0.7), Inches(2.4), Inches(5.6), Inches(4.3), [
    "User: \"Top 5 important bank emails this month\"",
    "detectIntent → important_emails",
    "Category detector → financial",
    "findImportantEmails(category=\"financial\")",
    "Runs ONE Gmail query (was 7 — 6x faster)",
    "Subject + from filters: payment, invoice, bill, statement, chase, amex, paypal, citi, …",
    "Score by recency, slice top N",
    "Return JSON → frontend renders Gmail-style card",
], size=12)

add_card(s, Inches(6.7), Inches(1.7), Inches(6.0), Inches(5.0))
add_text(s, Inches(6.9), Inches(1.85), Inches(5.6), Inches(0.5),
         "Categories supported", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(6.9), Inches(2.4), Inches(5.6), Inches(4.3), [
    "💰 Financial — banks, credit cards, payments, billing",
    "🏛 Government — IRS, USCIS, DMV, courts, .gov",
    "💼 Job — recruiters, interviews, offers, Workday",
    "⚠️ Urgent — deadlines, action-required, expires",
    "⭐ Starred — Gmail's user-flagged",
    "📄 Documents — contracts, signatures, DocuSign",
    "🌐 General — all 5 buckets unioned + ranked",
], size=13)

# ── Slide 7 — Evolution timeline ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Evolution", "From v1 \"can it speak\" to v16.6 \"can it actually help\"")
versions = [
    ("v1-v5", "Skeleton", "Chat → Groq → ElevenLabs. Could speak, couldn't do anything."),
    ("v6-v10", "Integrations", "Gmail OAuth + send + read. Calendar booking. Slack. Weather."),
    ("v11-v14", "UX polish", "Gmail-style email cards. Reply suggestions. Stop button."),
    ("v15.x", "Voice quality", "Echo blackout (1.5s), liveBar removal, deterministic summaries."),
    ("v16.0-v16.2", "Smart routing", "Multi-turn ack, bulk organize, importance ranker built."),
    ("v16.3", "Crash fix", "lastMsg const → let bug found and squashed."),
    ("v16.4", "Editable modal", "To/Subject/Message inputs in confirm-send dialog."),
    ("v16.5", "THE bug fix", "stripSystemPrefix bracket leak (the root cause of the email loop saga)."),
    ("v16.6", "Category ranker", "\"Bank emails\", \"job emails\" — single-bucket precision."),
]
y = Inches(1.75)
for i, (ver, name, desc) in enumerate(versions):
    row_y = y + Inches(i * 0.55)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), row_y, Inches(1.6), Inches(0.45))
    chip.adjustments[0] = 0.4; chip.fill.solid(); chip.fill.fore_color.rgb = TEAL
    chip.line.fill.background()
    add_text(s, Inches(0.5), row_y + Inches(0.05), Inches(1.6), Inches(0.4),
             ver, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.3), row_y + Inches(0.05), Inches(2.5), Inches(0.4),
             name, size=14, bold=True, color=ACCENT)
    add_text(s, Inches(5.0), row_y + Inches(0.05), Inches(8.0), Inches(0.4),
             desc, size=12, color=WHITE)

# ── Slide 8 — Live demo plan ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "Live demo", "What we'll see in the next 3 minutes")
demos = [
    ("1", "Greeting", "\"Hi\" → just a chat reply. No accidental email fetch.",
     "Proves the bracket-leak fix works."),
    ("2", "Smart category", "\"Top 5 important bank emails this month\"",
     "Importance ranker + financial filter."),
    ("3", "Sender search", "\"Show me my Workday emails\"",
     "Direct sender lookup."),
    ("4", "Bulk action", "\"Archive all from LinkedIn\"",
     "Bulk organize over up to 50 matches."),
    ("5", "Reply", "Tap \"Send This Reply\" on a job email; edit body in modal.",
     "Editable confirm-send."),
    ("6", "Wake word", "\"Hey ARIA, what's important today?\"",
     "Hands-free trigger + voice answer."),
]
for i, (num, what, sample, why) in enumerate(demos):
    row = i // 2; col = i % 2
    y = Inches(1.7 + row * 1.7); x = Inches(0.5 + col * 6.4)
    add_card(s, x, y, Inches(6.2), Inches(1.55))
    add_text(s, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.5),
             num, size=22, bold=True, color=TEAL)
    add_text(s, x + Inches(0.85), y + Inches(0.13), Inches(5.2), Inches(0.4),
             what, size=15, bold=True, color=WHITE)
    add_text(s, x + Inches(0.85), y + Inches(0.55), Inches(5.2), Inches(0.4),
             sample, size=12, color=ACCENT)
    add_text(s, x + Inches(0.85), y + Inches(1.0), Inches(5.2), Inches(0.5),
             why, size=11, color=DIM)

# ── Slide 9 — Numbers ──────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "By the numbers", "What it took to get here")
def stat(x, big, label, color):
    add_text(s, x, Inches(2.2), Inches(3), Inches(1.1),
             big, size=72, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.4), Inches(3), Inches(0.5),
             label, size=14, color=WHITE, align=PP_ALIGN.CENTER)
stat(Inches(0.4), "16.6", "current version", TEAL)
stat(Inches(3.5), "~2.2s", "voice round-trip", CYAN)
stat(Inches(6.7), "$0", "monthly cost (free tiers)", ACCENT)
stat(Inches(9.9), "6", "external integrations", TEAL)
add_card(s, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.0))
add_text(s, Inches(0.7), Inches(4.75), Inches(11.9), Inches(0.5),
         "Cost breakdown — fully usable on free tiers", size=16, bold=True, color=TEAL)
add_bullets(s, Inches(0.7), Inches(5.25), Inches(11.9), Inches(1.4), [
    "Cloudflare Workers — 100k requests / day free",
    "Groq Cloud — generous free tier, ~$0.59/M input tokens beyond",
    "ElevenLabs — 10k chars/month free (~5 minutes of speech)",
    "Gmail API + Google OAuth — free for personal use",
    "Supabase — 500MB free DB",
    "GitHub Pages — free hosting + free HTTPS",
], size=12)

# ── Slide 10 — What's next ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
slide_header(s, "What's next", "Where ARIA goes after the demo")
add_card(s, Inches(0.5), Inches(1.8), Inches(6.1), Inches(5.0))
add_text(s, Inches(0.7), Inches(1.95), Inches(5.7), Inches(0.5),
         "Near-term (weeks)", size=18, bold=True, color=TEAL)
add_bullets(s, Inches(0.7), Inches(2.5), Inches(5.7), Inches(4.2), [
    "Outlook + Microsoft Calendar parity",
    "Meeting summaries (Google Meet API)",
    "Smart unsubscribe (\"stop these newsletters\")",
    "Drafts before send (multi-turn editing)",
    "Mobile-first PWA install",
], size=14)
add_card(s, Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0))
add_text(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.5),
         "Longer-term (months)", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(4.2), [
    "Function-calling instead of regex routing",
    "Local Whisper STT (privacy-first option)",
    "Per-user fine-tuned importance signals",
    "Slack DMs, Notion notes, Linear tickets",
    "Persistent memory across sessions",
], size=14)

# ── Slide 11 — Thank you ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); add_bg(s)
add_text(s, Inches(0), Inches(2.5), Inches(13.33), Inches(1.0),
         "Thank you", size=72, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(3.7), Inches(13.33), Inches(0.5),
         "Questions?", size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(5.5), Inches(13.33), Inches(0.4),
         "ariaproxy.nihharksonee.workers.dev   ·   github.com/NKS-Coder/ARIA-voice-agent",
         size=14, color=DIM, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(6.1), Inches(13.33), Inches(0.4),
         "Nihhar K Sonee", size=14, color=ACCENT, align=PP_ALIGN.CENTER)

out = "/home/user/ARIA-voice-agent/ARIA-Demo.pptx"
prs.save(out)
print("✓ Wrote", out)
